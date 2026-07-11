#!/usr/bin/env python3
"""Generate 7 slide PNGs (1920x1080) + a matching PPTX for FrugalRouter."""

import os

from PIL import Image, ImageDraw, ImageFont

OUT = os.path.dirname(os.path.abspath(__file__))
W, H = 1920, 1080
BG = (14, 15, 18)
FG = (240, 240, 242)
DIM = (160, 163, 170)
RED = (237, 28, 36)          # AMD red accent
MONO_BG = (26, 28, 34)

FONT = "/System/Library/Fonts/HelveticaNeue.ttc"


def font(size, bold=False):
    # HelveticaNeue.ttc: index 0 regular, 1 medium/bold variants
    try:
        return ImageFont.truetype(FONT, size, index=1 if bold else 0)
    except OSError:
        return ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", size)


SLIDES = [
    {
        "kicker": "AMD DEVELOPER HACKATHON: ACT II — TRACK 1",
        "title": "FrugalRouter",
        "sub": "Hybrid Token-Efficient Routing Agent",
        "bullets": [],
        "footer": "ghcr.io/manan-tech/frugalrouter:v1   ·   github.com/manan-tech/frugalrouter",
    },
    {
        "kicker": "THE PROBLEM",
        "title": "Fewest tokens, above the accuracy gate",
        "bullets": [
            "19 hidden tasks across 8 capability categories",
            "Every Fireworks token counts against your score",
            "Grading box: 2 vCPU · 4 GB RAM · 10 minutes · no GPU",
        ],
    },
    {
        "kicker": "THE IDEA",
        "title": "Don't pick the cheapest cloud model.\nMake cloud calls unnecessary.",
        "bullets": [
            "Qwen3-1.7B + Qwen2.5-Coder-1.5B baked into the container",
            "llama.cpp, CPU-only, ~2.6 GB RAM — runs on the grading box itself",
            "Local inference costs exactly 0 tokens",
        ],
    },
    {
        "kicker": "THE CORE INNOVATION",
        "title": "Verification, not trust",
        "bullets": [
            "Math:  model writes Python — we EXECUTE it, majority-vote the number",
            "Code:  two independent implementations must AGREE on behavior",
            "Logic:  every assignment brute-forced, unique solution PROVEN",
            "Formats:  grammar-constrained decoding + programmatic checks",
        ],
    },
    {
        "kicker": "THE ROUTING LAYER",
        "title": "Escalate only what verification can't certify",
        "bullets": [
            "Confidence score per answer, from verification outcomes",
            "Weakest answers escalate cheapest-first to Fireworks AI",
            "Live-measured frugal champion: gpt-oss-120b @ low reasoning",
            "Hard global budget: 900 tokens  ·  budget 0 = pure zero-API mode",
        ],
    },
    {
        "kicker": "ENGINEERED FOR A HOSTILE CLOCK",
        "title": "The container cannot fail",
        "bullets": [
            "Startup tok/s probe drives a speed governor (full / lean / panic)",
            "results.json written atomically after every single task",
            "Watchdog flush at 8.5 min — TIMEOUT structurally impossible",
            "19 tasks in 232 s measured at grading constraints",
        ],
    },
    {
        "kicker": "RESULTS",
        "title": "Verified answers, almost free",
        "bullets": [
            "90% accuracy with ZERO Fireworks tokens (unseen 40-task suite)",
            "95% with escalation — expected spend ~300–700 tokens per run",
            "2.24 GB image · pure-stdlib Python · fully containerized",
            "llama.cpp · Qwen · Fireworks AI · Docker · GitHub Actions",
        ],
    },
]


def draw_slide(i, s):
    img = Image.new("RGB", (W, H), BG)
    d = ImageDraw.Draw(img)
    # accent bar + page dots
    d.rectangle([0, 0, 18, H], fill=RED)
    for j in range(len(SLIDES)):
        x = W - 60 - (len(SLIDES) - 1 - j) * 34
        d.ellipse([x, H - 70, x + 14, H - 56], fill=RED if j == i else (60, 62, 68))

    y = 130
    if s.get("kicker"):
        d.text((120, y), s["kicker"], font=font(34, bold=True), fill=RED)
        y += 90
    title_f = font(112 if i == 0 else 78, bold=True)
    for line in s["title"].split("\n"):
        d.text((120, y), line, font=title_f, fill=FG)
        y += (130 if i == 0 else 96)
    if s.get("sub"):
        d.text((120, y + 10), s["sub"], font=font(54), fill=DIM)
        y += 110
    y += 50
    for b in s["bullets"]:
        d.ellipse([126, y + 22, 146, y + 42], outline=RED, width=5)
        d.text((180, y), b, font=font(44), fill=FG)
        y += 92
    if s.get("footer"):
        d.rectangle([100, H - 190, W - 100, H - 110], fill=MONO_BG)
        d.text((130, H - 172), s["footer"], font=font(38), fill=DIM)
    img.save(f"{OUT}/slide{i + 1}.png")


for i, s in enumerate(SLIDES):
    draw_slide(i, s)
print(f"PNGs written to {OUT}")

# ---- PPTX ----
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.util import Emu, Pt  # noqa: E402

prs = Presentation()
prs.slide_width = Emu(12192000)   # 16:9
prs.slide_height = Emu(6858000)
blank = prs.slide_layouts[6]
for i in range(len(SLIDES)):
    slide = prs.slides.add_slide(blank)
    pic = slide.shapes.add_picture(f"{OUT}/slide{i + 1}.png", 0, 0,
                                   width=prs.slide_width, height=prs.slide_height)
prs.save(f"{OUT}/FrugalRouter_slides.pptx")
print("PPTX written")
